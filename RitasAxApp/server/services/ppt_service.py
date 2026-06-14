from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation
from pptx.util import Inches, Pt


def build_ppt(output_dir: Path, report_id: str, request_id: str, test_type: str, summary: Dict[str, Any], chart_paths: List[str], image_paths: List[str], llm_payload: str) -> Path:
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"[{test_type}] {request_id}"
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
    tf = textbox.text_frame
    tf.text = '요약'
    for k, v in summary.items():
        p = tf.add_paragraph()
        p.text = f"- {k}: {v}"
        p.font.size = Pt(16)

    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = 'Artifacts'
    y = 1.2
    for path in chart_paths[:3] + image_paths[:3]:
        p = slide2.shapes.add_textbox(Inches(0.7), Inches(y), Inches(8), Inches(0.4)).text_frame
        p.text = path
        y += 0.45

    out_path = output_dir / f"{test_type}_{request_id}_{report_id}.pptx"
    prs.save(out_path)
    return out_path
