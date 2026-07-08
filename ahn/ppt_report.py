"""PowerPoint report builder for AHN TEM/STEM/EDS/coating-layer analysis."""

from __future__ import annotations

import hashlib
import re
import tempfile
import xml.etree.ElementTree as ET
from copy import deepcopy
from pathlib import Path
from typing import Any, Iterable
from zipfile import ZipFile

from PIL import Image, ImageOps

from .docx_extract import extract_docx


def _require_pptx():
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.util import Inches, Pt
    except Exception as exc:  # pragma: no cover - depends on runtime packaging.
        raise RuntimeError(
            "AHN PPT 보고서를 생성하려면 python-pptx 패키지가 필요합니다."
        ) from exc
    return Presentation, RGBColor, PP_ALIGN, MSO_ANCHOR, Inches, Pt


Presentation, RGBColor, PP_ALIGN, MSO_ANCHOR, Inches, Pt = _require_pptx()

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
NAVY = RGBColor(31, 55, 87)
BLUE = RGBColor(47, 125, 211)
LIGHT_GRAY = RGBColor(245, 247, 250)
GRID_LINE = RGBColor(210, 220, 232)
TEXT = RGBColor(30, 47, 70)
TEMPLATE_PATH = Path(__file__).resolve().parent / "resources" / "templates" / "ahn_tem_template.pptx"
EMU_PER_INCH = 914400
PICTURE_SHAPE_TYPE = 13
TABLE_SHAPE_TYPE = 19
CAPTION_GAP = Inches(0.06)
CAPTION_HEIGHT = Inches(0.28)
CAPTION_FONT_SIZE = 9
XLSX_NS = {
    "x": "http://schemas.openxmlformats.org/spreadsheetml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}
TEMPLATE_SLIDES = {
    "tem": 0,
    "stem": 5,
    "stem_bf": 6,
    "eds_map": 7,
    "eds_line_first": 9,
    "eds_line_page": 10,
    "eds_table": 16,
    "coating_images": 19,
    "coating": 20,
}


def _chunks(items: list[Any], size: int) -> Iterable[list[Any]]:
    for index in range(0, len(items), size):
        yield items[index:index + size]


def _path(input_root: Path, relative: str) -> Path:
    return input_root / relative


def _new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _remove_shape(shape) -> None:
    shape.element.getparent().remove(shape.element)


def _clear_slides(prs) -> None:
    slide_id_list = prs.slides._sldIdLst
    for slide_id in list(slide_id_list):
        prs.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return " ".join(shape.text.split())


def _shape_bounds(shape) -> tuple[int, int, int, int]:
    return int(shape.left), int(shape.top), int(shape.width), int(shape.height)


def _slot_row_bucket(top: int) -> int:
    top_inches = top / EMU_PER_INCH
    return int((top_inches + 0.35) / 0.75)


def _sort_slots(slots: list[tuple[int, int, int, int]]) -> list[tuple[int, int, int, int]]:
    return sorted(slots, key=lambda slot: (_slot_row_bucket(slot[1]), slot[0]))


def _natural_key(value: str | Path) -> list[Any]:
    return [int(part) if part.isdigit() else part.lower() for part in re.split(r"(\d+)", str(value))]


def _magnification_sort_value(value: str) -> float:
    match = re.search(r"x?\s*(\d+(?:\.\d+)?)\s*([kKmM]?)", str(value or ""))
    if not match:
        return float("inf")
    scale = 1.0
    suffix = match.group(2).lower()
    if suffix == "k":
        scale = 1_000.0
    elif suffix == "m":
        scale = 1_000_000.0
    return float(match.group(1)) * scale


def _image_sort_key(item: dict[str, Any]) -> tuple[float, int, list[Any]]:
    sequence = item.get("sequence")
    sequence_value = sequence if isinstance(sequence, int) else 1_000_000_000
    return (
        _magnification_sort_value(str(item.get("magnification") or "")),
        sequence_value,
        _natural_key(str(item.get("file_name") or item.get("path") or "")),
    )


def _sort_image_items_for_report(image_items: list[dict[str, Any]]) -> list[dict[str, Any]]:
    return sorted(image_items, key=_image_sort_key)


def _fit_slot_to_canvas(slot: tuple[int, int, int, int], slide_width: int, slide_height: int) -> tuple[int, int, int, int]:
    left, top, width, height = slot
    margin = Inches(0.08)
    right = min(left + width, slide_width - margin)
    bottom = min(top + height, slide_height - margin)
    left = max(left, margin)
    top = max(top, margin)
    return left, top, max(1, right - left), max(1, bottom - top)


class AhnTemplate:
    def __init__(self, path: Path):
        self.path = path
        self.source = Presentation(path)
        self.output = Presentation(path)
        _strip_layout_footer(self.output)
        _clear_slides(self.output)

    def source_slide(self, key: str):
        return self.source.slides[TEMPLATE_SLIDES[key]]

    def picture_slots(self, key: str) -> list[tuple[int, int, int, int]]:
        return _sort_slots(
            [
                _fit_slot_to_canvas(_shape_bounds(shape), self.output.slide_width, self.output.slide_height)
                for shape in self.source_slide(key).shapes
                if shape.shape_type == PICTURE_SHAPE_TYPE
            ]
        )

    def label_slots(self, key: str) -> list[tuple[int, int, int, int]]:
        return _sort_slots(
            [
                _shape_bounds(shape)
                for shape in self.source_slide(key).shapes
                if getattr(shape, "has_text_frame", False)
                and _shape_text(shape).startswith("[")
                and _shape_text(shape).endswith("]")
            ]
        )

    def table_slots(self, key: str) -> list[tuple[int, int, int, int]]:
        return _sort_slots(
            [
                _shape_bounds(shape)
                for shape in self.source_slide(key).shapes
                if shape.shape_type == TABLE_SHAPE_TYPE
            ]
        )

    def new_slide(self, key: str, title: str):
        slide = self.output.slides.add_slide(self.output.slide_layouts[0])
        for shape in list(slide.shapes):
            _remove_shape(shape)

        for shape in self.source_slide(key).shapes:
            slide.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")

        _strip_template_payload(slide, key)
        _replace_template_title(slide, title)
        return slide


def _strip_layout_footer(prs) -> None:
    for layout in prs.slide_layouts:
        for shape in list(layout.shapes):
            text = _shape_text(shape)
            if "‹#›" in text or "<#>" in text or ("/ 22" in text and shape.top > Inches(6.5)):
                _remove_shape(shape)


def _strip_template_payload(slide, key: str) -> None:
    for shape in list(slide.shapes):
        text = _shape_text(shape)
        if shape.shape_type == PICTURE_SHAPE_TYPE:
            _remove_shape(shape)
            continue
        if getattr(shape, "has_text_frame", False) and text.startswith("[") and text.endswith("]"):
            _remove_shape(shape)
            continue
        if shape.shape_type == TABLE_SHAPE_TYPE:
            if key == "coating" and shape.left < Inches(7.8):
                continue
            _remove_shape(shape)


def _replace_template_title(slide, title: str) -> None:
    title_text = f"■ {title}"
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if _shape_text(shape).startswith("■"):
            shape.left = Inches(0.2)
            shape.width = Inches(10.35)
            frame = shape.text_frame
            frame.clear()
            para = frame.paragraphs[0]
            run = para.add_run()
            run.text = title_text
            _set_run(run, size=16, bold=True, color=TEXT)
            return
    _add_text(slide, title_text, Inches(0.09), Inches(0.76), Inches(10.66), Inches(0.45), size=16, bold=True)


def _set_run(run, *, size: int, bold: bool = False, color=TEXT) -> None:
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_text(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    *,
    size: int = 14,
    bold: bool = False,
    color=TEXT,
    align=None,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    para = frame.paragraphs[0]
    if align is not None:
        para.alignment = align
    run = para.add_run()
    run.text = text
    _set_run(run, size=size, bold=bold, color=color)
    return box


def _add_header(slide, section_title: str) -> None:
    _add_text(
        slide,
        "TEM 분석 결과",
        Inches(10.35),
        Inches(0.18),
        Inches(2.5),
        Inches(0.35),
        size=15,
        bold=True,
        color=NAVY,
    )
    _add_text(
        slide,
        f"■ {section_title}",
        Inches(0.35),
        Inches(0.42),
        Inches(9.6),
        Inches(0.48),
        size=20,
        bold=True,
        color=NAVY,
    )
    line = slide.shapes.add_shape(1, Inches(0.35), Inches(1.02), Inches(12.55), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.color.rgb = BLUE


def _convert_image(path: Path, tmp_dir: Path) -> Path:
    cache_dir = tmp_dir / "converted-images"
    cache_dir.mkdir(parents=True, exist_ok=True)
    digest = hashlib.sha1(str(path.resolve()).encode("utf-8")).hexdigest()[:16]
    target = cache_dir / f"{digest}.jpg"
    if target.exists():
        return target
    with Image.open(path) as image:
        image = ImageOps.exif_transpose(image)
        if image.mode == "RGBA":
            background = Image.new("RGB", image.size, (255, 255, 255))
            background.paste(image, mask=image.getchannel("A"))
            image = background
        elif image.mode != "RGB":
            image = image.convert("RGB")
        image.thumbnail((1600, 1600), Image.Resampling.LANCZOS)
        image.save(target, format="JPEG", quality=88, optimize=True)
    return target


def _image_size(path: Path) -> tuple[int, int]:
    with Image.open(path) as image:
        return image.size


def _image_aspect(path: Path) -> float:
    width, height = _image_size(path)
    if height <= 0:
        return 0.0
    return width / height


def _add_fit_picture(slide, path: Path, left, top, width, height, tmp_dir: Path):
    converted = _convert_image(path, tmp_dir)
    image_w, image_h = _image_size(converted)
    if image_w <= 0 or image_h <= 0:
        return None
    scale = min(width / image_w, height / image_h)
    pic_w = int(image_w * scale)
    pic_h = int(image_h * scale)
    pic_left = int(left + (width - pic_w) / 2)
    pic_top = int(top + (height - pic_h) / 2)
    return slide.shapes.add_picture(str(converted), pic_left, pic_top, width=pic_w, height=pic_h)


def _caption_slot_for_image(slot: tuple[int, int, int, int]) -> tuple[int, int, int, int]:
    return (slot[0], slot[1] + slot[3] + CAPTION_GAP, slot[2], CAPTION_HEIGHT)


def _add_label(slide, label: str, slot: tuple[int, int, int, int]) -> None:
    _add_text(
        slide,
        label,
        slot[0],
        slot[1],
        slot[2],
        slot[3],
        size=CAPTION_FONT_SIZE,
        bold=True,
        color=TEXT,
        align=PP_ALIGN.CENTER,
    )


def _add_template_image_items(
    slide,
    image_items: list[dict[str, Any]],
    image_slots: list[tuple[int, int, int, int]],
    label_slots: list[tuple[int, int, int, int]],
    input_root: Path,
    tmp_dir: Path,
    *,
    show_labels: bool = True,
) -> None:
    for index, item in enumerate(image_items[: len(image_slots)]):
        slot = image_slots[index]
        _add_fit_picture(slide, _path(input_root, item["path"]), *slot, tmp_dir)
        if show_labels:
            label = item.get("magnification") or Path(item["file_name"]).stem
            if label and not label.startswith("["):
                label = f"[{label}]"
            _add_label(slide, label, _caption_slot_for_image(slot))


def _add_template_paths(
    slide,
    paths: list[Path],
    image_slots: list[tuple[int, int, int, int]],
    tmp_dir: Path,
) -> None:
    for index, path in enumerate(paths[: len(image_slots)]):
        _add_fit_picture(slide, path, *image_slots[index], tmp_dir)


def _add_image_grid(
    slide,
    image_items: list[dict[str, Any]],
    input_root: Path,
    tmp_dir: Path,
    *,
    left,
    top,
    width,
    height,
    cols: int,
    rows: int,
    show_labels: bool = True,
) -> None:
    gap = Inches(0.1)
    label_h = CAPTION_HEIGHT if show_labels else 0
    cell_w = (width - gap * (cols - 1)) / cols
    cell_h = (height - gap * (rows - 1)) / rows
    for index, item in enumerate(image_items[: cols * rows]):
        col = index % cols
        row = index // cols
        cell_left = left + col * (cell_w + gap)
        cell_top = top + row * (cell_h + gap)
        image_h = cell_h - label_h
        _add_fit_picture(slide, _path(input_root, item["path"]), cell_left, cell_top, cell_w, image_h, tmp_dir)
        if not show_labels:
            continue
        label = item.get("magnification") or Path(item["file_name"]).stem
        if label and not label.startswith("["):
            label = f"[{label}]"
        _add_text(
            slide,
            label,
            cell_left,
            cell_top + image_h,
            cell_w,
            label_h,
            size=CAPTION_FONT_SIZE,
            bold=True,
            color=TEXT,
            align=PP_ALIGN.CENTER,
        )


def _add_image_grid_slides(
    prs,
    template: AhnTemplate | None,
    template_key: str,
    title: str,
    image_items: list[dict[str, Any]],
    input_root: Path,
    tmp_dir: Path,
    *,
    per_slide: int = 8,
) -> None:
    if not image_items:
        return
    sorted_items = _sort_image_items_for_report(image_items)
    for page_index, chunk in enumerate(_chunks(sorted_items, per_slide), start=1):
        suffix = f" ({page_index})" if len(image_items) > per_slide else ""
        if template:
            slide = template.new_slide(template_key, title + suffix)
            _add_image_grid(
                slide,
                chunk,
                input_root,
                tmp_dir,
                left=Inches(0.36),
                top=Inches(1.25),
                width=Inches(10.45),
                height=Inches(5.85),
                cols=4,
                rows=2,
            )
        else:
            slide = _new_slide(prs)
            _add_header(slide, title + suffix)
            _add_image_grid(
                slide,
                chunk,
                input_root,
                tmp_dir,
                left=Inches(0.65),
                top=Inches(1.25),
                width=Inches(12.0),
                height=Inches(5.95),
                cols=4,
                rows=2,
            )


def _add_table(
    slide,
    rows: list[list[str]],
    left,
    top,
    width,
    height,
    *,
    font_size: int = 9,
    column_widths: list[int] | None = None,
    alignments: list[Any] | None = None,
    draw_borders: bool = False,
):
    if not rows:
        return None
    cols = max(len(row) for row in rows)
    table_shape = slide.shapes.add_table(len(rows), cols, left, top, width, height)
    table = table_shape.table
    if column_widths:
        for col_index, col_width in enumerate(column_widths[:cols]):
            table.columns[col_index].width = int(col_width)
    for row_index, row in enumerate(rows):
        for col_index in range(cols):
            cell = table.cell(row_index, col_index)
            value = row[col_index] if col_index < len(row) else ""
            cell.text = str(value)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.03)
            cell.margin_right = Inches(0.03)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.text_frame.word_wrap = True
            alignment = (
                alignments[col_index]
                if alignments and col_index < len(alignments)
                else PP_ALIGN.CENTER
            )
            if row_index == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
                text_color = RGBColor(255, 255, 255)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                text_color = TEXT
            for para in cell.text_frame.paragraphs:
                para.alignment = alignment
                if not para.runs:
                    para.add_run()
                for run in para.runs:
                    _set_run(run, size=font_size, bold=row_index == 0, color=text_color)
            if draw_borders:
                _set_cell_borders(cell)
    return table_shape


def _set_cell_borders(cell, *, color: str = "1F3757", width: int = 12700) -> None:
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement

    tc_pr = cell._tc.get_or_add_tcPr()
    for edge in ("lnL", "lnR", "lnT", "lnB"):
        existing = tc_pr.find(qn(f"a:{edge}"))
        if existing is not None:
            tc_pr.remove(existing)
        line = OxmlElement(f"a:{edge}")
        line.set("w", str(width))
        line.set("cap", "flat")
        line.set("cmpd", "sng")
        line.set("algn", "ctr")

        solid_fill = OxmlElement("a:solidFill")
        rgb = OxmlElement("a:srgbClr")
        rgb.set("val", color)
        solid_fill.append(rgb)
        line.append(solid_fill)

        dash = OxmlElement("a:prstDash")
        dash.set("val", "solid")
        line.append(dash)
        tc_pr.append(line)


def _xlsx_shared_strings(zip_file: ZipFile) -> list[str]:
    try:
        root = ET.fromstring(zip_file.read("xl/sharedStrings.xml"))
    except KeyError:
        return []
    return [
        "".join(node.text or "" for node in item.findall(".//x:t", XLSX_NS)).strip()
        for item in root.findall("x:si", XLSX_NS)
    ]


def _xlsx_sheet_paths(zip_file: ZipFile) -> list[tuple[str, str]]:
    try:
        workbook = ET.fromstring(zip_file.read("xl/workbook.xml"))
        relationships = ET.fromstring(zip_file.read("xl/_rels/workbook.xml.rels"))
    except KeyError:
        return []
    targets = {
        relationship.attrib.get("Id", ""): relationship.attrib.get("Target", "")
        for relationship in relationships.findall("rel:Relationship", XLSX_NS)
    }
    paths: list[tuple[str, str]] = []
    for sheet in workbook.findall(".//x:sheet", XLSX_NS):
        relationship_id = sheet.attrib.get(f"{{{XLSX_NS['r']}}}id", "")
        target = targets.get(relationship_id, "")
        if not target:
            continue
        if target.startswith("/"):
            path = target.lstrip("/")
        elif target.startswith("xl/"):
            path = target
        else:
            path = f"xl/{target}"
        paths.append((sheet.attrib.get("name", ""), path))
    return paths


def _xlsx_col_index(cell_ref: str) -> int:
    letters = "".join(ch for ch in cell_ref if ch.isalpha()).upper()
    value = 0
    for letter in letters:
        value = value * 26 + ord(letter) - 64
    return max(0, value - 1)


def _format_xlsx_text(value: str) -> str:
    text = value.strip()
    if not text:
        return ""
    try:
        numeric = float(text)
    except ValueError:
        return text
    if numeric.is_integer():
        return str(int(numeric))
    if abs(numeric) >= 1000 or 0 < abs(numeric) < 0.001:
        return f"{numeric:.3g}"
    return f"{numeric:.3f}".rstrip("0").rstrip(".")


def _xlsx_cell_text(cell: ET.Element, shared_strings: list[str]) -> str:
    if cell.attrib.get("t") == "inlineStr":
        return "".join(node.text or "" for node in cell.findall(".//x:t", XLSX_NS)).strip()
    value = cell.find("x:v", XLSX_NS)
    if value is None or value.text is None:
        return ""
    text = value.text.strip()
    if cell.attrib.get("t") == "s":
        try:
            return shared_strings[int(text)]
        except (IndexError, ValueError):
            return text
    return _format_xlsx_text(text)


def _read_xlsx_tables(path: Path, *, max_tables: int = 2, max_rows: int = 12, max_cols: int = 10) -> list[list[list[str]]]:
    tables: list[list[list[str]]] = []
    try:
        with ZipFile(path) as zip_file:
            shared_strings = _xlsx_shared_strings(zip_file)
            for _sheet_name, sheet_path in _xlsx_sheet_paths(zip_file):
                try:
                    sheet = ET.fromstring(zip_file.read(sheet_path))
                except KeyError:
                    continue
                rows: list[list[str]] = []
                for row in sheet.findall(".//x:sheetData/x:row", XLSX_NS):
                    values: list[str] = []
                    for cell in row.findall("x:c", XLSX_NS):
                        col_index = _xlsx_col_index(cell.attrib.get("r", "A1"))
                        if col_index >= max_cols:
                            continue
                        while len(values) <= col_index:
                            values.append("")
                        values[col_index] = _xlsx_cell_text(cell, shared_strings)
                    while values and not values[-1]:
                        values.pop()
                    if any(values):
                        rows.append(values)
                    if len(rows) >= max_rows:
                        break
                if len(rows) >= 2 and max(len(row) for row in rows) >= 2:
                    tables.append(rows)
                if len(tables) >= max_tables:
                    break
    except (OSError, ET.ParseError):
        return []
    return tables


def _normalize_lookup(value: str) -> str:
    return "".join(ch.lower() for ch in value if ch.isalnum())


def _matching_spreadsheet_tables(data: dict[str, Any], report: dict[str, Any], input_root: Path) -> list[list[list[str]]]:
    sample_key = _normalize_lookup(str(report.get("sample_name") or ""))
    title_key = _normalize_lookup(str(report.get("title") or ""))
    candidates: list[Path] = []
    fallback: list[Path] = []
    for item in data.get("spreadsheets") or []:
        path = _path(input_root, item.get("path", ""))
        if path.suffix.lower() != ".xlsx":
            continue
        lookup = _normalize_lookup(path.stem)
        fallback.append(path)
        if (sample_key and sample_key in lookup) or (title_key and title_key in lookup):
            candidates.append(path)
    if not candidates and len(fallback) == 1:
        candidates = fallback
    for path in candidates:
        tables = _read_xlsx_tables(path)
        if tables:
            return tables
    return []


def _eds_report_tables(data: dict[str, Any], report: dict[str, Any], input_root: Path, docx_tables: list[list[list[str]]]) -> list[list[list[str]]]:
    if docx_tables:
        return docx_tables
    return _matching_spreadsheet_tables(data, report, input_root)


def _build_tem(prs, template: AhnTemplate | None, data: dict[str, Any], input_root: Path, tmp_dir: Path) -> None:
    for sample in data.get("tem_samples") or []:
        _add_image_grid_slides(
            prs,
            template,
            "tem",
            f"TEM 이미지 분석결과 : [{sample['sample_name']}]",
            sample.get("images") or [],
            input_root,
            tmp_dir,
        )


def _build_stem(prs, template: AhnTemplate | None, data: dict[str, Any], input_root: Path, tmp_dir: Path) -> None:
    for sample in data.get("stem_samples") or []:
        if sample.get("images"):
            _add_image_grid_slides(
                prs,
                template,
                "stem",
                f"STEM 이미지 분석결과 : [{sample['sample_name']}]",
                sample.get("images") or [],
                input_root,
                tmp_dir,
            )
        if sample.get("bf_images"):
            _add_image_grid_slides(
                prs,
                template,
                "stem_bf",
                f"STEM BF 이미지 분석결과 : [{sample['sample_name']}]",
                sample.get("bf_images") or [],
                input_root,
                tmp_dir,
            )


def _eds_anchor_slot() -> tuple[int, int, int, int]:
    return Inches(0.2), Inches(1.25), Inches(4.4), Inches(3.64)


def _eds_right_slot() -> tuple[int, int, int, int]:
    return Inches(4.25), Inches(1.32), Inches(6.45), Inches(4.48)


def _eds_full_grid_slot() -> tuple[int, int, int, int]:
    return Inches(0.26), Inches(1.38), Inches(10.5), Inches(5.55)


def _eds_line_anchor_slot() -> tuple[int, int, int, int]:
    return Inches(0.2), Inches(1.25), Inches(4.35), Inches(3.44)


def _eds_line_top_slot() -> tuple[int, int, int, int]:
    return Inches(4.7), Inches(1.36), Inches(6.1), Inches(2.21)


def _eds_line_bottom_slot() -> tuple[int, int, int, int]:
    return Inches(4.7), Inches(3.66), Inches(6.1), Inches(2.94)


def _eds_table_slots(count: int, area: tuple[int, int, int, int]) -> list[tuple[int, int, int, int]]:
    left, top, width, height = area
    if count <= 1:
        return [(left, top, width, height)]
    gap = Inches(0.16)
    slot_h = int((height - gap) / 2)
    return [
        (left, top, width, slot_h),
        (left, top + slot_h + gap, width, height - slot_h - gap),
    ]


def _add_eds_anchor_image(slide, first_image: Path | None, tmp_dir: Path) -> None:
    if first_image:
        _add_fit_picture(slide, first_image, *_eds_anchor_slot(), tmp_dir)


def _add_eds_right_grid(
    slide,
    images: list[Path],
    tmp_dir: Path,
    *,
    cols: int,
    rows: int,
) -> None:
    if not images:
        return
    _add_absolute_image_grid(
        slide,
        [{"path": str(path), "file_name": path.name, "magnification": ""} for path in images],
        tmp_dir,
        *_eds_right_slot(),
        cols,
        rows,
    )


def _add_eds_anchor_grid(
    slide,
    first_image: Path | None,
    images: list[Path],
    tmp_dir: Path,
    *,
    cols: int,
    rows: int,
) -> None:
    _add_eds_anchor_image(slide, first_image, tmp_dir)
    _add_eds_right_grid(slide, images, tmp_dir, cols=cols, rows=rows)


def _add_eds_line_overview(
    slide,
    first_image: Path | None,
    line_images: list[Path],
    tmp_dir: Path,
) -> None:
    if first_image:
        _add_fit_picture(slide, first_image, *_eds_line_anchor_slot(), tmp_dir)
    if line_images[:1]:
        _add_fit_picture(slide, line_images[0], *_eds_line_top_slot(), tmp_dir)
    if line_images[1:2]:
        _add_fit_picture(slide, line_images[1], *_eds_line_bottom_slot(), tmp_dir)


def _add_eds_slide(prs, template: AhnTemplate | None, key: str, title: str):
    if template:
        return template.new_slide(key, f"STEM EDS 분석결과 : [{title}]")
    slide = _new_slide(prs)
    _add_header(slide, f"STEM EDS 분석결과 : [{title}]")
    return slide


def _add_eds_map_slide(
    prs,
    template: AhnTemplate | None,
    title: str,
    first_image: Path | None,
    map_images: list[Path],
    tmp_dir: Path,
) -> None:
    if not first_image and not map_images:
        return
    slide = _add_eds_slide(prs, template, "eds_map", title)
    _add_eds_anchor_grid(slide, first_image, map_images, tmp_dir, cols=3, rows=2)


def _add_eds_map_pages(
    prs,
    template: AhnTemplate | None,
    title: str,
    images: list[Path],
    tmp_dir: Path,
) -> None:
    if not images:
        return
    first_image = images[0]
    map_images = images[1:]
    if not map_images:
        _add_eds_map_slide(prs, template, title, first_image, [], tmp_dir)
        return
    for page, chunk in enumerate(_chunks(map_images, 6), start=1):
        page_title = title if page == 1 else f"{title}_Data{page}"
        _add_eds_map_slide(prs, template, page_title, first_image, chunk, tmp_dir)


def _add_absolute_image_grid(slide, image_items: list[dict[str, Any]], tmp_dir: Path, left, top, width, height, cols: int, rows: int) -> None:
    gap = Inches(0.035)
    cell_w = (width - gap * (cols - 1)) / cols
    cell_h = (height - gap * (rows - 1)) / rows
    for index, item in enumerate(image_items[: cols * rows]):
        col = index % cols
        row = index // cols
        _add_fit_picture(
            slide,
            Path(item["path"]),
            left + col * (cell_w + gap),
            top + row * (cell_h + gap),
            cell_w,
            cell_h,
            tmp_dir,
        )


def _add_eds_image_pages(
    prs,
    template: AhnTemplate | None,
    title: str,
    images: list[Path],
    tmp_dir: Path,
    *,
    start_page: int = 1,
    template_key: str = "eds_line_page",
    cols: int = 3,
    rows: int = 2,
    append_data_suffix: bool = True,
) -> None:
    per_page = max(1, cols * rows)
    for offset, chunk in enumerate(_chunks(images, per_page)):
        page = start_page + offset
        if append_data_suffix:
            page_title = f"{title}_Data{page}"
        else:
            page_title = title if offset == 0 else f"{title}_Graph{offset + 1}"
        slide = _add_eds_slide(prs, template, template_key, page_title)
        _add_absolute_image_grid(
            slide,
            [{"path": str(path), "file_name": path.name, "magnification": ""} for path in chunk],
            tmp_dir,
            *_eds_full_grid_slot(),
            cols,
            rows,
        )


def _line_eds_group_starts(images: list[Path]) -> list[int]:
    starts: list[int] = []
    for index, image in enumerate(images):
        if index + 2 >= len(images):
            continue
        try:
            aspect = _image_aspect(image)
            next_aspect = _image_aspect(images[index + 1])
            third_aspect = _image_aspect(images[index + 2])
        except Exception:
            continue
        if 1.1 <= aspect <= 2.2 and next_aspect >= 2.4 and third_aspect >= 2.4:
            starts.append(index)
    if not starts:
        return [0] if images else []
    if starts[0] != 0:
        starts.insert(0, 0)
    return starts


def _split_line_eds_groups(images: list[Path]) -> list[list[Path]]:
    starts = _line_eds_group_starts(images)
    groups: list[list[Path]] = []
    for index, start in enumerate(starts):
        end = starts[index + 1] if index + 1 < len(starts) else len(images)
        group = images[start:end]
        if group:
            groups.append(group)
    return groups


def _add_eds_line_group(
    prs,
    template: AhnTemplate | None,
    title: str,
    group: list[Path],
    tmp_dir: Path,
    *,
    data_index: int,
) -> None:
    if not group:
        return
    page_title = f"{title}_Data{data_index}"
    slide = _add_eds_slide(prs, template, "eds_line_first", page_title)
    _add_eds_line_overview(slide, group[0], group[1:3], tmp_dir)
    _add_eds_image_pages(
        prs,
        template,
        page_title,
        group[3:],
        tmp_dir,
        start_page=1,
        template_key="eds_line_page",
        cols=2,
        rows=3,
        append_data_suffix=False,
    )


def _add_eds_line_pages(
    prs,
    template: AhnTemplate | None,
    title: str,
    images: list[Path],
    tmp_dir: Path,
) -> None:
    for data_index, group in enumerate(_split_line_eds_groups(images), start=1):
        _add_eds_line_group(prs, template, title, group, tmp_dir, data_index=data_index)


def _add_eds_anchor_grid_pages(
    prs,
    template: AhnTemplate | None,
    title: str,
    first_image: Path | None,
    images: list[Path],
    tmp_dir: Path,
    *,
    start_page: int,
) -> None:
    """Add EDS continuation slides while keeping the first Word image anchored left."""
    if not images:
        return
    for page, chunk in enumerate(_chunks(images, 6), start=start_page):
        if template:
            slide = template.new_slide("eds_line_first", f"STEM EDS 분석결과 : [{title}_Data{page}]")
        else:
            slide = _new_slide(prs)
            _add_header(slide, f"STEM EDS 분석결과 : [{title}_Data{page}]")
        _add_eds_anchor_grid(slide, first_image, chunk, tmp_dir, cols=2, rows=3)


def _is_point_spectrum_table(table: list[list[str]]) -> bool:
    return bool(
        table
        and table[0]
        and "spectrum" in str(table[0][0]).strip().lower()
        and len(table) >= 2
    )


def _point_table_unit(table: list[list[str]]) -> str:
    if not table or not table[0]:
        return ""
    unit = str(table[0][-1]).strip()
    return unit if unit.lower() in {"at%", "wt%"} else ""


def _point_table_value_width(table: list[list[str]]) -> int:
    if not table or not table[0]:
        return 0
    header = [str(value).strip() for value in table[0]]
    width = len(header)
    if width >= 2 and header[-2].lower() == "total" and header[-1].lower() in {"at%", "wt%"}:
        return width - 2
    if width >= 1 and header[-1].lower() == "total":
        return width - 1
    return width


def _normalized_point_composition_rows(
    table: list[list[str]],
    *,
    body_index: int | None = None,
) -> list[list[str]]:
    if not _is_point_spectrum_table(table):
        return []
    width = _point_table_value_width(table)
    if width <= 0:
        return []
    unit = _point_table_unit(table)
    header = [str(value) for value in table[0][:width]]
    source_rows = table[1:]
    if body_index is not None:
        if body_index < 0 or body_index >= len(source_rows):
            return []
        source_rows = source_rows[body_index:body_index + 1]
    rows = [[unit or "Composition"] + [""] * (width - 1), header]
    for row in source_rows:
        rows.append([str(row[index]) if index < len(row) else "" for index in range(width)])
    return rows


def _point_table_font_size(rows: list[list[str]]) -> int:
    cols = max((len(row) for row in rows), default=0)
    if cols > 11 or len(rows) > 9:
        return 5
    if cols > 8 or len(rows) > 6:
        return 6
    return 7


def _restyle_merged_header(cell, text: str, font_size: int) -> None:
    cell.text = text
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        if not paragraph.runs:
            paragraph.add_run()
        for run in paragraph.runs:
            _set_run(run, size=font_size, bold=True, color=RGBColor(255, 255, 255))


def _add_point_composition_table(
    slide,
    table: list[list[str]],
    slot: tuple[int, int, int, int],
    *,
    body_index: int | None = None,
) -> Any | None:
    rows = _normalized_point_composition_rows(table, body_index=body_index)
    if not rows:
        return None
    font_size = _point_table_font_size(rows)
    table_shape = _add_table(
        slide,
        rows,
        *slot,
        font_size=font_size,
        alignments=[PP_ALIGN.CENTER] * max(len(row) for row in rows),
        draw_borders=True,
    )
    if table_shape is None:
        return None
    cols = len(rows[0])
    if cols > 1:
        header_cell = table_shape.table.cell(0, 0)
        header_cell.merge(table_shape.table.cell(0, cols - 1))
        _restyle_merged_header(header_cell, rows[0][0], font_size)
    return table_shape


def _valid_point_tables(tables: list[list[list[str]]]) -> list[list[list[str]]]:
    return [table for table in tables[:2] if _is_point_spectrum_table(table)]


def _eds_point_anchor_slot() -> tuple[int, int, int, int]:
    return Inches(0.26), Inches(1.28), Inches(4.25), Inches(4.44)


def _eds_point_table_area() -> tuple[int, int, int, int]:
    return Inches(4.65), Inches(1.28), Inches(6.08), Inches(2.02)


def _eds_point_graph_slot() -> tuple[int, int, int, int]:
    return Inches(4.65), Inches(3.47), Inches(6.08), Inches(3.02)


def _add_eds_point_summary_slide(
    prs,
    template: AhnTemplate | None,
    title: str,
    first_image: Path | None,
    point_tables: list[list[list[str]]],
    tmp_dir: Path,
) -> None:
    slide = _add_eds_slide(prs, template, "eds_table", title)
    _add_eds_anchor_image(slide, first_image, tmp_dir)
    slots = _eds_table_slots(len(point_tables), _eds_right_slot())
    for table, slot in zip(point_tables, slots):
        _add_point_composition_table(slide, table, slot)


def _point_spectrum_labels(point_tables: list[list[list[str]]]) -> list[str]:
    if not point_tables:
        return []
    return [
        str(row[0]).strip() or f"Spectrum {index + 1}"
        for index, row in enumerate(point_tables[0][1:])
    ]


def _add_eds_point_detail_slide(
    prs,
    template: AhnTemplate | None,
    title: str,
    first_image: Path | None,
    spectrum_image: Path | None,
    point_tables: list[list[list[str]]],
    body_index: int,
    spectrum_label: str,
    tmp_dir: Path,
) -> None:
    slide = _add_eds_slide(prs, template, "eds_line_first", f"{title}_{spectrum_label}")
    if first_image:
        _add_fit_picture(slide, first_image, *_eds_point_anchor_slot(), tmp_dir)
    slots = _eds_table_slots(len(point_tables), _eds_point_table_area())
    for table, slot in zip(point_tables, slots):
        _add_point_composition_table(slide, table, slot, body_index=body_index)
    if spectrum_image:
        _add_fit_picture(slide, spectrum_image, *_eds_point_graph_slot(), tmp_dir)


def _add_eds_point_pages(
    prs,
    template: AhnTemplate | None,
    title: str,
    images: list[Path],
    tables: list[list[list[str]]],
    tmp_dir: Path,
) -> None:
    point_tables = _valid_point_tables(tables)
    if not point_tables:
        _add_eds_tables_slide(prs, template, title, images[:1], tables, tmp_dir)
        _add_eds_anchor_grid_pages(
            prs,
            template,
            title,
            images[0] if images else None,
            images[1:],
            tmp_dir,
            start_page=1,
        )
        return

    first_image = images[0] if images else None
    spectrum_images = images[1:]
    _add_eds_point_summary_slide(prs, template, title, first_image, point_tables, tmp_dir)
    labels = _point_spectrum_labels(point_tables)
    for index, label in enumerate(labels):
        _add_eds_point_detail_slide(
            prs,
            template,
            title,
            first_image,
            spectrum_images[index] if index < len(spectrum_images) else None,
            point_tables,
            index,
            label,
            tmp_dir,
        )
    if len(spectrum_images) > len(labels):
        _add_eds_image_pages(
            prs,
            template,
            title,
            spectrum_images[len(labels):],
            tmp_dir,
            start_page=len(labels) + 1,
            template_key="eds_line_page",
            cols=2,
            rows=3,
        )


def _add_eds_tables_slide(prs, template: AhnTemplate | None, title: str, images: list[Path], tables: list[list[list[str]]], tmp_dir: Path) -> None:
    if not images and not tables:
        return
    if template:
        slide = template.new_slide("eds_table", f"STEM EDS 분석결과 : [{title}]")
    else:
        slide = _new_slide(prs)
        _add_header(slide, f"STEM EDS 분석결과 : [{title}]")
    _add_eds_anchor_image(slide, images[0] if images else None, tmp_dir)
    if not tables:
        return
    table_area = _eds_right_slot() if images else _eds_full_grid_slot()
    slots = _eds_table_slots(len(tables[:2]), table_area)
    for table, slot in zip(tables[:2], slots):
        normalized = [row[:10] for row in table[:14]]
        _add_table(slide, normalized, *slot, font_size=7)


def _build_eds(prs, template: AhnTemplate | None, data: dict[str, Any], input_root: Path, tmp_dir: Path) -> None:
    for report in data.get("eds_reports") or []:
        docx_path = _path(input_root, report["path"])
        extract_dir = tmp_dir / "docx" / docx_path.stem
        extracted = extract_docx(docx_path, extract_dir)
        images = extracted.media_paths
        tables = _eds_report_tables(data, report, input_root, extracted.tables)
        analysis_type = str(report.get("analysis_type") or "").upper()
        title = report.get("title") or docx_path.stem
        if analysis_type == "MAP":
            _add_eds_map_pages(prs, template, title, images, tmp_dir)
        elif analysis_type == "LINE":
            _add_eds_line_pages(prs, template, title, images, tmp_dir)
        elif analysis_type == "POINT":
            _add_eds_point_pages(prs, template, title, images, tables, tmp_dir)
        else:
            _add_eds_map_pages(prs, template, title, images, tmp_dir)


def _coating_grid_shape(count: int) -> tuple[int, int, int]:
    if count <= 9:
        return 3, 3, 9
    if count <= 12:
        return 4, 3, 12
    return 4, 4, 16


def _format_nm(value: Any) -> str:
    if value is None:
        return "검토 필요"
    try:
        return f"{float(value):.2f}"
    except (TypeError, ValueError):
        return "검토 필요"


def _coating_note(value: Any) -> str:
    note = str(value or "")
    return (
        note.replace("OCR 라벨", "라벨")
        .replace("OCR 후보값", "후보값")
        .replace("OCR 실패", "자동 판독 실패")
        .replace("OCR 엔진", "자동 판독 엔진")
    )


def _coating_rows(measurements: list[dict[str, Any]]) -> list[list[str]]:
    rows = [["측정개소", "두께(nm)"]]
    values: list[float] = []
    row_number = 1
    for item in measurements:
        item_values = item.get("thickness_values_nm") or []
        if not item_values and item.get("thickness_nm") is not None:
            item_values = [item.get("thickness_nm")]
        if not item_values:
            rows.append([str(row_number), "검토 필요"])
            row_number += 1
            continue
        for value in item_values:
            try:
                numeric_value = float(value)
            except (TypeError, ValueError):
                rows.append([str(row_number), "검토 필요"])
                row_number += 1
                continue
            values.append(numeric_value)
            rows.append([str(row_number), _format_nm(numeric_value)])
            row_number += 1
    average = sum(values) / len(values) if values else None
    rows.append(["전체 평균", _format_nm(average)])
    return rows


def _coating_column_widths(width: int) -> list[int]:
    first = int(width * 0.42)
    return [first, int(width - first)]


def _coating_table_font_size(row_count: int) -> int:
    if row_count > 32:
        return 9
    if row_count > 24:
        return 10
    return 12


def _remove_tables(slide) -> None:
    for shape in list(slide.shapes):
        if shape.shape_type == TABLE_SHAPE_TYPE:
            _remove_shape(shape)


def _add_coating_table_slide(
    prs,
    template: AhnTemplate | None,
    title: str,
    measurements: list[dict[str, Any]],
) -> None:
    table_rows = _coating_rows(measurements)
    if template:
        slide = template.new_slide("coating", f"{title} 두께 요약")
        _remove_tables(slide)
        left, top, width = Inches(1.65), Inches(1.35), Inches(7.55)
        height = min(Inches(5.75), max(Inches(1.0), Inches(0.33 * len(table_rows))))
    else:
        slide = _new_slide(prs)
        _add_header(slide, f"{title} 두께 요약")
        left, top, width = Inches(2.2), Inches(1.35), Inches(8.9)
        height = min(Inches(5.75), max(Inches(1.0), Inches(0.33 * len(table_rows))))
    _add_table(
        slide,
        table_rows,
        left,
        top,
        width,
        height,
        font_size=_coating_table_font_size(len(table_rows)),
        column_widths=_coating_column_widths(width),
        alignments=[PP_ALIGN.CENTER, PP_ALIGN.CENTER],
        draw_borders=True,
    )


def _add_coating_image_only_slides(
    prs,
    template: AhnTemplate | None,
    title: str,
    measurements: list[dict[str, Any]],
    input_root: Path,
    tmp_dir: Path,
) -> None:
    if not measurements:
        return
    if template:
        slots = template.picture_slots("coating_images")
        per_slide = max(1, len(slots))
    else:
        per_slide = 12
    for page_index, chunk in enumerate(_chunks(measurements, per_slide), start=1):
        suffix = f" ({page_index})" if len(measurements) > per_slide else ""
        if template:
            slide = template.new_slide("coating_images", title + suffix)
            image_items = [
                {
                    "path": item["path"],
                    "file_name": item["file_name"],
                    "magnification": item.get("magnification") or "",
                }
                for item in chunk
            ]
            _add_template_image_items(
                slide,
                image_items,
                template.picture_slots("coating_images"),
                [],
                input_root,
                tmp_dir,
                show_labels=False,
            )
        else:
            slide = _new_slide(prs)
            _add_header(slide, title + suffix)
            image_items = [
                {
                    "path": item["path"],
                    "file_name": item["file_name"],
                    "magnification": item.get("magnification") or "",
                }
                for item in chunk
            ]
            _add_image_grid(
                slide,
                image_items,
                input_root,
                tmp_dir,
                left=Inches(0.65),
                top=Inches(1.25),
                width=Inches(9.8),
                height=Inches(5.95),
                cols=4,
                rows=3,
                show_labels=False,
            )


def _coating_final_grid_shape(count: int) -> tuple[int, int]:
    if count <= 2:
        return max(1, count), 1
    if count <= 4:
        return 2, 2
    if count <= 6:
        return 3, 2
    if count <= 9:
        return 3, 3
    return 4, 3


def _add_coating_images_with_table_slide(
    prs,
    template: AhnTemplate | None,
    title: str,
    image_measurements: list[dict[str, Any]],
    all_measurements: list[dict[str, Any]],
    input_root: Path,
    tmp_dir: Path,
) -> None:
    if template:
        slide = template.new_slide("coating", title)
        _remove_tables(slide)
    else:
        slide = _new_slide(prs)
        _add_header(slide, title)
    image_items = [
        {
            "path": item["path"],
            "file_name": item["file_name"],
            "magnification": item.get("magnification") or "",
        }
        for item in image_measurements
    ]
    cols, rows = _coating_final_grid_shape(len(image_items))
    _add_image_grid(
        slide,
        image_items,
        input_root,
        tmp_dir,
        left=Inches(0.72),
        top=Inches(1.22),
        width=Inches(6.72),
        height=Inches(5.95),
        cols=cols,
        rows=rows,
        show_labels=False,
    )
    table_rows = _coating_rows(all_measurements)
    _add_table(
        slide,
        table_rows,
        Inches(7.75),
        Inches(1.2),
        Inches(2.85),
        Inches(5.95),
        font_size=_coating_table_font_size(len(table_rows)),
        column_widths=_coating_column_widths(Inches(2.85)),
        alignments=[PP_ALIGN.CENTER, PP_ALIGN.CENTER],
        draw_borders=True,
    )


def _build_coating(prs, template: AhnTemplate | None, data: dict[str, Any], input_root: Path, tmp_dir: Path) -> None:
    for sample in data.get("coating_samples") or []:
        measurements = sample.get("measurements") or []
        if not measurements:
            continue
        base_title = f"TEM 코팅층 두께 분석 결과 : [{sample['sample_name']}]"
        if len(measurements) >= 10:
            image_page_size = 12
            remainder = len(measurements) % image_page_size
            final_size = remainder or min(image_page_size, len(measurements))
            image_only_measurements = measurements[:-final_size]
            final_measurements = measurements[-final_size:]
            if image_only_measurements:
                _add_coating_image_only_slides(
                    prs,
                    template,
                    base_title,
                    image_only_measurements,
                    input_root,
                    tmp_dir,
                )
            final_page = (len(image_only_measurements) // image_page_size) + 1
            final_title = (
                f"{base_title} ({final_page})"
                if image_only_measurements
                else base_title
            )
            _add_coating_images_with_table_slide(
                prs,
                template,
                final_title,
                final_measurements,
                measurements,
                input_root,
                tmp_dir,
            )
            continue
        if template:
            per_slide = max(1, len(template.picture_slots("coating")))
        else:
            cols, rows, per_slide = _coating_grid_shape(len(measurements))
        for page_index, chunk in enumerate(_chunks(measurements, per_slide), start=1):
            suffix = f" ({page_index})" if len(measurements) > per_slide else ""
            title = f"{base_title}{suffix}"
            if template:
                slide = template.new_slide("coating", title)
            else:
                slide = _new_slide(prs)
                _add_header(slide, title)
            image_items = [
                {
                    "path": item["path"],
                    "file_name": item["file_name"],
                    "magnification": item.get("magnification") or "",
                }
                for item in chunk
            ]
            if template:
                _add_template_image_items(
                    slide,
                    image_items,
                    template.picture_slots("coating"),
                    [],
                    input_root,
                    tmp_dir,
                    show_labels=False,
                )
                table_rows = _coating_rows(chunk)
                table_slots = template.table_slots("coating")
                table_slot = table_slots[1] if len(table_slots) > 1 else (Inches(8.01), Inches(1.21), Inches(2.36), Inches(5.77))
                _add_table(
                    slide,
                    table_rows,
                    *table_slot,
                    font_size=_coating_table_font_size(len(table_rows)),
                    column_widths=_coating_column_widths(table_slot[2]),
                    alignments=[PP_ALIGN.CENTER, PP_ALIGN.CENTER],
                    draw_borders=True,
                )
            else:
                _add_image_grid(
                    slide,
                    image_items,
                    input_root,
                    tmp_dir,
                    left=Inches(0.55),
                    top=Inches(1.35),
                    width=Inches(8.15),
                    height=Inches(5.75),
                    cols=cols,
                    rows=rows,
                    show_labels=False,
                )
                table_rows = _coating_rows(chunk)
                table_height = min(Inches(5.75), Inches(0.32 * len(table_rows)))
                _add_table(
                    slide,
                    table_rows,
                    Inches(9.0),
                    Inches(1.35),
                    Inches(3.55),
                    table_height,
                    font_size=_coating_table_font_size(len(table_rows)),
                    column_widths=_coating_column_widths(Inches(3.55)),
                    alignments=[PP_ALIGN.CENTER, PP_ALIGN.CENTER],
                    draw_borders=True,
                )


def build_pptx(project: Any, output_path: str | Path) -> Path:
    """Build an AHN PowerPoint report from collected project data."""
    data = project.to_dict() if hasattr(project, "to_dict") else dict(project)
    input_root = Path(data["input_root"])
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    template = AhnTemplate(TEMPLATE_PATH) if TEMPLATE_PATH.exists() else None
    if template:
        prs = template.output
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

    with tempfile.TemporaryDirectory(prefix="ahn-pptx-") as tmp:
        tmp_dir = Path(tmp)
        _build_tem(prs, template, data, input_root, tmp_dir)
        _build_stem(prs, template, data, input_root, tmp_dir)
        _build_eds(prs, template, data, input_root, tmp_dir)
        _build_coating(prs, template, data, input_root, tmp_dir)
        if not prs.slides:
            if template:
                slide = template.new_slide("tem", "AHN 분석결과")
            else:
                slide = _new_slide(prs)
                _add_header(slide, "AHN 분석결과")
            _add_text(slide, "입력 폴더에서 보고서 생성 대상 데이터를 찾지 못했습니다.", Inches(0.7), Inches(1.6), Inches(8), Inches(0.5))
        prs.save(output)
    return output
