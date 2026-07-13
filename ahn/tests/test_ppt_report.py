from __future__ import annotations

from pathlib import Path
from types import SimpleNamespace
from zipfile import ZipFile

import pytest
from PIL import Image

pptx = pytest.importorskip("pptx")
Presentation = pptx.Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ahn.ppt_report import (
    _coating_rows,
    _coating_table_font_size,
    _point_table_column_widths,
    build_pptx,
)


def _write_image(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", (120, 90), (220, 224, 230)).save(path)


def _write_sized_image(path: Path, size: tuple[int, int]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", size, (220, 224, 230)).save(path)


def _write_minimal_xlsx(path: Path, rows: list[list[str]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    sheet_rows = []
    for row_index, row in enumerate(rows, start=1):
        cells = []
        for col_index, value in enumerate(row):
            col = chr(ord("A") + col_index)
            cells.append(
                f'<c r="{col}{row_index}" t="inlineStr"><is><t>{value}</t></is></c>'
            )
        sheet_rows.append(f'<row r="{row_index}">{"".join(cells)}</row>')
    with ZipFile(path, "w") as zip_file:
        zip_file.writestr(
            "xl/workbook.xml",
            (
                '<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" '
                'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
                '<sheets><sheet name="line 1" sheetId="1" r:id="rId1"/></sheets></workbook>'
            ),
        )
        zip_file.writestr(
            "xl/_rels/workbook.xml.rels",
            (
                '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" '
                'Target="worksheets/sheet1.xml"/></Relationships>'
            ),
        )
        zip_file.writestr(
            "xl/worksheets/sheet1.xml",
            (
                '<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
                f'<sheetData>{"".join(sheet_rows)}</sheetData></worksheet>'
            ),
        )


def _pptx_table_text(path: Path) -> str:
    prs = Presentation(path)
    values: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    values.extend(cell.text for cell in row.cells)
    return "\n".join(values)


def _pptx_picture_counts(path: Path) -> list[int]:
    prs = Presentation(path)
    return [sum(1 for shape in slide.shapes if shape.shape_type == 13) for slide in prs.slides]


def _black_grid_rule_count(slide) -> int:
    count = 0
    for shape in slide.shapes:
        try:
            color = str(shape.fill.fore_color.rgb)
        except Exception:
            continue
        if color == "000000" and (shape.width <= Inches(0.02) or shape.height <= Inches(0.02)):
            count += 1
    return count


def _pictures(slide):
    return sorted(
        [shape for shape in slide.shapes if shape.shape_type == 13],
        key=lambda shape: (shape.left, shape.top),
    )


def _caption_texts(slide) -> list[str]:
    captions = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text.strip()
        if text.startswith("[") and text.endswith("]"):
            captions.append((shape.top, shape.left, text))
    return [text for _top, _left, text in sorted(captions)]


def test_coating_table_rows_hide_ocr_word_and_keep_readable_font_policy() -> None:
    rows = _coating_rows(
        [
            {
                "index": 1,
                "thickness_values_nm": [2.0, 3.0],
                "note": "OCR 라벨 2개 추출",
                "ocr_warnings": ["OCR 후보값 없음"],
            }
        ]
    )

    joined = "\n".join("\t".join(row) for row in rows)
    assert rows[0] == ["측정개소", "두께(nm)"]
    assert "OCR" not in joined
    assert "1\t2.00" in joined
    assert "2\t3.00" in joined
    assert _coating_table_font_size(20) == 12


def _base_project(root: Path) -> dict:
    return {
        "experiment": "AHN-TEM",
        "input_root": str(root),
        "generated_at": "2026-07-08T00:00:00+00:00",
        "folders": {},
        "tem_samples": [],
        "stem_samples": [],
        "eds_reports": [],
        "spreadsheets": [],
        "coating_samples": [],
        "summary": {},
    }


def _project_with_sections(root: Path, sections: set[str]) -> dict:
    image_path = root / "image.png"
    _write_image(image_path)
    report_path = root / "report" / "report.docx"
    report_path.parent.mkdir(parents=True, exist_ok=True)
    report_path.write_bytes(b"placeholder")

    image_record = {
        "path": image_path.name,
        "file_name": image_path.name,
        "sample_name": "S1",
        "magnification": "x100k",
        "sequence": 1,
        "kind": "",
    }
    project = _base_project(root)
    if "tem" in sections:
        project["tem_samples"] = [{"sample_name": "TEM-A", "images": [image_record]}]
    if "stem" in sections:
        project["stem_samples"] = [{"sample_name": "STEM-A", "images": [image_record], "bf_images": []}]
    if "eds" in sections:
        project["eds_reports"] = [
            {
                "path": "report/report.docx",
                "file_name": "report.docx",
                "title": "EDS-A",
                "sample_name": "EDS-A",
                "analysis_type": "MAP",
            }
        ]
    if "coating" in sections:
        project["coating_samples"] = [
            {
                "sample_name": "Scale-A",
                "measurements": [
                    {
                        "index": 1,
                        "path": image_path.name,
                        "file_name": image_path.name,
                        "magnification": "x100k",
                        "thickness_nm": 2.5,
                        "thickness_values_nm": [2.0, 3.0],
                        "ocr_text": "",
                        "note": "라벨 2개 추출",
                        "ocr_review_required": False,
                        "ocr_warnings": [],
                    }
                ],
            }
        ]
    return project


@pytest.mark.parametrize(
    ("sections", "expected_titles"),
    [
        ({"tem", "stem", "eds", "coating"}, ["■ TEM 이미지 분석결과", "■ STEM 이미지 분석결과", "■ STEM EDS 분석결과", "■ TEM 코팅층"]),
        ({"tem", "eds", "coating"}, ["■ TEM 이미지 분석결과", "■ STEM EDS 분석결과", "■ TEM 코팅층"]),
        ({"stem", "coating"}, ["■ STEM 이미지 분석결과", "■ TEM 코팅층"]),
        ({"eds"}, ["■ STEM EDS 분석결과"]),
    ],
)
def test_build_pptx_includes_only_available_sections(tmp_path, monkeypatch, sections, expected_titles) -> None:
    image_path = tmp_path / "image.png"

    def fake_extract_docx(_docx_path: Path, _extract_dir: Path):
        return SimpleNamespace(media_paths=[image_path], tables=[])

    monkeypatch.setattr("ahn.ppt_report.extract_docx", fake_extract_docx)
    project = _project_with_sections(tmp_path, set(sections))
    output = tmp_path / "report.pptx"

    build_pptx(project, output)

    prs = Presentation(output)
    slide_texts = [" ".join(shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)) for slide in prs.slides]
    all_text = "\n".join(slide_texts)
    assert len(prs.slides) == len(expected_titles)
    for title in expected_titles:
        assert title in all_text
    absent_titles = {
        "tem": "■ TEM 이미지 분석결과",
        "stem": "■ STEM 이미지 분석결과",
        "eds": "■ STEM EDS 분석결과",
        "coating": "■ TEM 코팅층",
    }
    for section, title in absent_titles.items():
        if section not in sections:
            assert title not in all_text


def test_point_eds_docx_tables_are_included(tmp_path, monkeypatch) -> None:
    image_path = tmp_path / "image.png"
    _write_image(image_path)
    report_path = tmp_path / "report" / "line.docx"
    report_path.parent.mkdir(parents=True, exist_ok=True)
    report_path.write_bytes(b"placeholder")

    def fake_extract_docx(_docx_path: Path, _extract_dir: Path):
        return SimpleNamespace(
            media_paths=[image_path],
            tables=[[["Element", "Wt%"], ["C", "12.3"], ["O", "45.6"]]],
        )

    monkeypatch.setattr("ahn.ppt_report.extract_docx", fake_extract_docx)
    project = _base_project(tmp_path)
    project["eds_reports"] = [
        {
            "path": "report/line.docx",
            "file_name": "line.docx",
            "title": "0283 point",
            "sample_name": "0283",
            "analysis_type": "POINT",
        }
    ]
    output = tmp_path / "report.pptx"

    build_pptx(project, output)

    table_text = _pptx_table_text(output)
    assert "Element" in table_text
    assert "12.3" in table_text


def test_point_eds_uses_matching_spreadsheet_when_docx_has_no_tables(tmp_path, monkeypatch) -> None:
    image_path = tmp_path / "image.png"
    _write_image(image_path)
    report_path = tmp_path / "report" / "line.docx"
    xlsx_path = tmp_path / "report" / "0283 line scan raw data.xlsx"
    report_path.parent.mkdir(parents=True, exist_ok=True)
    report_path.write_bytes(b"placeholder")
    _write_minimal_xlsx(
        xlsx_path,
        [["Point", "Distance (um)", "C Wt%"], ["1", "0", "62.9"], ["2", "0.012", "52.8"]],
    )

    def fake_extract_docx(_docx_path: Path, _extract_dir: Path):
        return SimpleNamespace(media_paths=[image_path], tables=[])

    monkeypatch.setattr("ahn.ppt_report.extract_docx", fake_extract_docx)
    project = _base_project(tmp_path)
    project["eds_reports"] = [
        {
            "path": "report/line.docx",
            "file_name": "line.docx",
            "title": "0283 point",
            "sample_name": "0283",
            "analysis_type": "POINT",
        }
    ]
    project["spreadsheets"] = [
        {"path": "report/0283 line scan raw data.xlsx", "file_name": "0283 line scan raw data.xlsx"}
    ]
    output = tmp_path / "report.pptx"

    build_pptx(project, output)

    table_text = _pptx_table_text(output)
    assert "Distance (um)" in table_text
    assert "62.9" in table_text


def test_line_eds_continuation_pages_use_full_width_graph_grid(tmp_path, monkeypatch) -> None:
    images = []
    for index in range(10):
        image_path = tmp_path / f"eds-{index}.png"
        image_path.parent.mkdir(parents=True, exist_ok=True)
        Image.new("RGB", (240, 100), (220, 224, 230)).save(image_path)
        images.append(image_path)
    report_path = tmp_path / "report" / "line.docx"
    report_path.parent.mkdir(parents=True, exist_ok=True)
    report_path.write_bytes(b"placeholder")

    def fake_extract_docx(_docx_path: Path, _extract_dir: Path):
        return SimpleNamespace(media_paths=images, tables=[])

    monkeypatch.setattr("ahn.ppt_report.extract_docx", fake_extract_docx)
    project = _base_project(tmp_path)
    project["eds_reports"] = [
        {
            "path": "report/line.docx",
            "file_name": "line.docx",
            "title": "0283 line scan",
            "sample_name": "0283",
            "analysis_type": "LINE",
        }
    ]
    output = tmp_path / "report.pptx"

    build_pptx(project, output)

    assert _pptx_picture_counts(output) == [3, 6, 1]
    prs = Presentation(output)
    first_slide_pictures = _pictures(prs.slides[0])
    anchor = first_slide_pictures[0]
    assert anchor.left <= Inches(0.3)
    assert anchor.width >= Inches(4.0)
    right_pictures = [shape for shape in first_slide_pictures if shape.left >= Inches(4.5)]
    assert len(right_pictures) == 2
    assert all(shape.width >= Inches(5.0) for shape in right_pictures)
    continuation_pictures = _pictures(prs.slides[1])
    assert len(continuation_pictures) == 6
    left_groups = {round(shape.left / Inches(1), 1) for shape in continuation_pictures}
    top_groups = {round(shape.top / Inches(1), 1) for shape in continuation_pictures}
    assert len(left_groups) == 2
    assert len(top_groups) == 3
    assert all(shape.width >= Inches(4.0) for shape in continuation_pictures)


def test_line_eds_repeats_overview_and_graph_grid_per_data_block(tmp_path, monkeypatch) -> None:
    pattern = [(840, 559), (863, 212), (857, 281)] + [(802, 252)] * 6
    images = []
    for index, size in enumerate(pattern * 2):
        image_path = tmp_path / f"line-block-{index}.png"
        _write_sized_image(image_path, size)
        images.append(image_path)
    report_path = tmp_path / "report" / "line.docx"
    report_path.parent.mkdir(parents=True, exist_ok=True)
    report_path.write_bytes(b"placeholder")

    def fake_extract_docx(_docx_path: Path, _extract_dir: Path):
        return SimpleNamespace(media_paths=images, tables=[])

    monkeypatch.setattr("ahn.ppt_report.extract_docx", fake_extract_docx)
    project = _base_project(tmp_path)
    project["eds_reports"] = [
        {
            "path": "report/line.docx",
            "file_name": "line.docx",
            "title": "0817 line scan",
            "sample_name": "0817",
            "analysis_type": "LINE",
        }
    ]
    output = tmp_path / "report.pptx"

    build_pptx(project, output)

    assert _pptx_picture_counts(output) == [3, 6, 3, 6]
    prs = Presentation(output)
    titles = [" ".join(shape.text.split()) for slide in prs.slides for shape in slide.shapes if getattr(shape, "has_text_frame", False)]
    assert any("0817 line scan_Data1" in title for title in titles)
    assert any("0817 line scan_Data2" in title for title in titles)
    assert not any("Data1_Data" in title or "Data2_Data" in title for title in titles)
    assert len(_pictures(prs.slides[1])) == 6
    assert all(shape.left >= Inches(0.2) for shape in _pictures(prs.slides[1]))
    assert len(_pictures(prs.slides[3])) == 6
    assert all(shape.left >= Inches(0.2) for shape in _pictures(prs.slides[3]))


def test_map_eds_pages_reuse_first_image_as_left_anchor(tmp_path, monkeypatch) -> None:
    images = []
    for index in range(8):
        image_path = tmp_path / f"map-{index}.png"
        _write_image(image_path)
        images.append(image_path)
    report_path = tmp_path / "report" / "map.docx"
    report_path.parent.mkdir(parents=True, exist_ok=True)
    report_path.write_bytes(b"placeholder")

    def fake_extract_docx(_docx_path: Path, _extract_dir: Path):
        return SimpleNamespace(media_paths=images, tables=[])

    monkeypatch.setattr("ahn.ppt_report.extract_docx", fake_extract_docx)
    project = _base_project(tmp_path)
    project["eds_reports"] = [
        {
            "path": "report/map.docx",
            "file_name": "map.docx",
            "title": "0817 MAP",
            "sample_name": "0817",
            "analysis_type": "MAP",
        }
    ]
    output = tmp_path / "report.pptx"

    build_pptx(project, output)

    assert _pptx_picture_counts(output) == [7, 2]
    prs = Presentation(output)
    for slide in prs.slides:
        pictures = _pictures(slide)
        anchor = pictures[0]
        assert anchor.left <= Inches(0.3)
        assert anchor.width >= Inches(4.0)
    assert len([shape for shape in _pictures(prs.slides[0]) if shape.left >= Inches(4.0)]) == 6
    first_page_right = [shape for shape in _pictures(prs.slides[0]) if shape.left >= Inches(4.0)]
    assert len({round(shape.left / Inches(1), 2) for shape in first_page_right}) == 2
    assert len({round(shape.top / Inches(1), 2) for shape in first_page_right}) == 3

    single_right_picture = [shape for shape in _pictures(prs.slides[1]) if shape.left >= Inches(4.0)]
    assert len(single_right_picture) == 1
    assert single_right_picture[0].left == min(shape.left for shape in first_page_right)
    assert single_right_picture[0].top == min(shape.top for shape in first_page_right)
    assert single_right_picture[0].width == first_page_right[0].width
    assert single_right_picture[0].height == first_page_right[0].height


def test_map_eds_partial_page_keeps_fixed_two_by_three_slots(tmp_path, monkeypatch) -> None:
    images = []
    for index in range(9):
        image_path = tmp_path / f"map-partial-{index}.png"
        _write_image(image_path)
        images.append(image_path)
    report_path = tmp_path / "report" / "map-partial.docx"
    report_path.parent.mkdir(parents=True, exist_ok=True)
    report_path.write_bytes(b"placeholder")

    monkeypatch.setattr(
        "ahn.ppt_report.extract_docx",
        lambda _docx_path, _extract_dir: SimpleNamespace(media_paths=images, tables=[]),
    )
    project = _base_project(tmp_path)
    project["eds_reports"] = [
        {
            "path": "report/map-partial.docx",
            "file_name": "map-partial.docx",
            "title": "Sample 2 MAP2",
            "sample_name": "Sample 2",
            "analysis_type": "MAP",
        }
    ]
    output = tmp_path / "map-partial.pptx"

    build_pptx(project, output)

    prs = Presentation(output)
    full_page = [shape for shape in _pictures(prs.slides[0]) if shape.left >= Inches(4.0)]
    partial_page = [shape for shape in _pictures(prs.slides[1]) if shape.left >= Inches(4.0)]
    assert len(partial_page) == 2
    assert partial_page[0].top == partial_page[1].top == min(shape.top for shape in full_page)
    assert [shape.left for shape in partial_page] == sorted({shape.left for shape in full_page})
    assert all(shape.width == full_page[0].width for shape in partial_page)
    assert all(shape.height == full_page[0].height for shape in partial_page)


def test_point_eds_detail_pages_keep_first_image_on_left(tmp_path, monkeypatch) -> None:
    images = []
    for index in range(4):
        image_path = tmp_path / f"point-{index}.png"
        _write_image(image_path)
        images.append(image_path)
    report_path = tmp_path / "report" / "point.docx"
    report_path.parent.mkdir(parents=True, exist_ok=True)
    report_path.write_bytes(b"placeholder")

    def fake_extract_docx(_docx_path: Path, _extract_dir: Path):
        return SimpleNamespace(
            media_paths=images,
            tables=[[["Element", "Wt%"], ["C", "12.3"]]],
        )

    monkeypatch.setattr("ahn.ppt_report.extract_docx", fake_extract_docx)
    project = _base_project(tmp_path)
    project["eds_reports"] = [
        {
            "path": "report/point.docx",
            "file_name": "point.docx",
            "title": "0283 point",
            "sample_name": "0283",
            "analysis_type": "POINT",
        }
    ]
    output = tmp_path / "report.pptx"

    build_pptx(project, output)

    assert _pptx_picture_counts(output) == [1, 4]
    prs = Presentation(output)
    anchor = _pictures(prs.slides[0])[0]
    assert anchor.left <= Inches(0.3)
    assert anchor.width >= Inches(4.3)
    table_lefts = [
        shape.left
        for shape in prs.slides[0].shapes
        if getattr(shape, "has_table", False)
    ]
    assert table_lefts
    assert min(table_lefts) - (anchor.left + anchor.width) <= Inches(0.12)
    for shape in prs.slides[0].shapes:
        if not getattr(shape, "has_table", False):
            continue
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    assert paragraph.alignment == PP_ALIGN.CENTER


def test_point_eds_spectrum_tables_create_row_detail_slides(tmp_path, monkeypatch) -> None:
    images = []
    for index in range(3):
        image_path = tmp_path / f"point-spectrum-{index}.png"
        _write_sized_image(image_path, (820, 548))
        images.append(image_path)
    report_path = tmp_path / "reports" / "point.docx"
    report_path.parent.mkdir(parents=True, exist_ok=True)
    report_path.write_bytes(b"placeholder")

    at_table = [
        ["Spectrum Label", "C", "N", "O", "Al", "Si", "Cr", "Mn", "Fe", "Cu", "Mo", "Total", "At%"],
        ["Spectrum 34", "16.38", "0.00", "27.73", "0.15", "6.59", "0.16", "0.19", "48.49", "0.05", "0.26", "100.00", "Project 1/0647 Point2"],
        ["Spectrum 35", "30.57", "0.00", "25.36", "0.10", "12.26", "0.37", "0.00", "31.13", "0.21", "0.00", "100.00", "Project 1/0647 Point2"],
    ]
    wt_table = [
        ["Spectrum Label", "C", "N", "O", "Al", "Si", "Cr", "Mn", "Fe", "Cu", "Mo", "Total", "Wt%"],
        ["Spectrum 34", "5.49", "0.00", "12.38", "0.11", "5.16", "0.23", "0.29", "75.56", "0.10", "0.68", "100.00", "Project 1/0647 Point2"],
        ["Spectrum 35", "12.70", "0.00", "14.04", "0.09", "11.91", "0.66", "0.00", "60.15", "0.45", "0.00", "100.00", "Project 1/0647 Point2"],
    ]

    def fake_extract_docx(_docx_path: Path, _extract_dir: Path):
        return SimpleNamespace(media_paths=images, tables=[at_table, wt_table])

    monkeypatch.setattr("ahn.ppt_report.extract_docx", fake_extract_docx)
    project = _base_project(tmp_path)
    project["eds_reports"] = [
        {
            "path": "reports/point.docx",
            "file_name": "point.docx",
            "title": "0647 Point2",
            "sample_name": "0647",
            "analysis_type": "POINT",
        }
    ]
    output = tmp_path / "report.pptx"

    build_pptx(project, output)

    prs = Presentation(output)
    assert len(prs.slides) == 4
    assert _pptx_picture_counts(output) == [1, 1, 2, 2]
    table_counts = [
        sum(1 for shape in slide.shapes if getattr(shape, "has_table", False))
        for slide in prs.slides
    ]
    assert table_counts == [1, 1, 2, 2]
    summary_units = []
    for slide in list(prs.slides)[:2]:
        tables = [shape.table for shape in slide.shapes if getattr(shape, "has_table", False)]
        assert len(tables) == 1
        summary_units.append(tables[0].cell(0, 0).text)
        anchor = _pictures(slide)[0]
        assert anchor.left <= Inches(0.3)
        assert anchor.width >= Inches(4.2)
    assert summary_units == ["Wt%", "At%"]
    summary_table_fonts = []
    for slide in list(prs.slides)[:2]:
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            assert shape.table.columns[0].width >= Inches(0.8)
            assert shape.table.columns[0].width > shape.table.columns[1].width
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.size:
                                summary_table_fonts.append(run.font.size)
    assert summary_table_fonts
    assert min(summary_table_fonts) >= Pt(9)
    detail_table_fonts = []
    for shape in prs.slides[2].shapes:
        if not getattr(shape, "has_table", False):
            continue
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size:
                            detail_table_fonts.append(run.font.size)
    assert detail_table_fonts
    assert min(detail_table_fonts) >= Pt(9)
    table_text = _pptx_table_text(output)
    assert "At%" in table_text
    assert "Wt%" in table_text
    assert "Spectrum 34" in table_text
    assert "Spectrum 35" in table_text
    assert "Total" not in table_text
    titles = [
        " ".join(shape.text.split())
        for slide in prs.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    ]
    assert any("0647 Point2_Spectrum 34" in title for title in titles)
    assert any("0647 Point2_Spectrum 35" in title for title in titles)


def test_point_table_widths_keep_spectrum_label_wide_and_other_columns_even() -> None:
    widths = _point_table_column_widths(
        Inches(5.0),
        [["Spectrum Label", "C", "O", "Si", "Cr", "Fe", "Cu", "Mo"]],
    )

    assert widths is not None
    assert widths[0] >= Inches(1.2)
    assert widths[0] > widths[1]
    assert max(widths[1:]) - min(widths[1:]) <= 1
    assert sum(widths) == Inches(5.0)


def test_large_coating_sample_uses_image_pages_then_summary_table(tmp_path) -> None:
    project = _base_project(tmp_path)
    measurements = []
    for index in range(1, 15):
        image_path = tmp_path / f"scale-{index}.png"
        _write_image(image_path)
        measurements.append(
            {
                "index": index,
                "path": image_path.name,
                "file_name": image_path.name,
                "magnification": "",
                "thickness_nm": float(index),
                "thickness_values_nm": [float(index)],
            }
        )
    project["coating_samples"] = [{"sample_name": "Scale-A", "measurements": measurements}]
    output = tmp_path / "report.pptx"

    build_pptx(project, output)

    prs = Presentation(output)
    assert len(prs.slides) == 2
    table_counts = [sum(1 for shape in slide.shapes if getattr(shape, "has_table", False)) for slide in prs.slides]
    assert table_counts[0] == 0
    assert table_counts[1] == 0
    assert _pptx_picture_counts(output) == [12, 3]
    pictures = _pictures(prs.slides[1])
    table_picture = pictures[-1]
    assert table_picture.width <= Inches(3.0)
    assert table_picture.left >= Inches(7.6)
    assert _black_grid_rule_count(prs.slides[1]) == 0


def test_grid_slides_place_remainder_images_from_upper_left_and_sort_by_magnification(tmp_path) -> None:
    project = _base_project(tmp_path)
    records = []
    for file_name in [
        "HR Camera_3_600kX.png",
        "HR Camera_1_20kX.png",
        "HR Camera_2_100kX.png",
        "HR Camera_4_300kX.png",
        "HR Camera_5_300kX.png",
        "HR Camera_6_300kX.png",
        "HR Camera_7_300kX.png",
        "HR Camera_8_300kX.png",
        "HR Camera_9_300kX.png",
    ]:
        image_path = tmp_path / file_name
        _write_image(image_path)
        mag = file_name.split("_")[2].replace("X.png", "")
        records.append(
            {
                "path": file_name,
                "file_name": file_name,
                "sample_name": "STEM-A",
                "magnification": f"x{mag}",
                "sequence": int(file_name.split("_")[1]),
                "kind": "STEM",
            }
        )
    project["stem_samples"] = [{"sample_name": "STEM-A", "images": records, "bf_images": []}]
    output = tmp_path / "report.pptx"

    build_pptx(project, output)

    prs = Presentation(output)
    assert _caption_texts(prs.slides[0])[:3] == ["[x20k]", "[x100k]", "[x300k]"]
    remainder_picture = _pictures(prs.slides[1])[0]
    assert remainder_picture.left <= Inches(0.5)
    assert remainder_picture.top <= Inches(1.7)
